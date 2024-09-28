import os
import json
import time
import shutil
from pathlib import Path
from fastapi import FastAPI, UploadFile, File, Form, HTTPException
from fastapi.responses import FileResponse
from fastapi.middleware.cors import CORSMiddleware
from pdf2docx import Converter
from pptx import Presentation
from docx import Document
import fitz  # PyMuPDF for extracting text from PDFs

app = FastAPI()

# Allow CORS for your frontend on Vercel
app.add_middleware(
    CORSMiddleware,
    allow_origins=["https://my-html-gray.vercel.app"],  # Change this to your Vercel frontend URL for production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

CONVERTED_FOLDER = Path('converted')
UPLOADED_FOLDER = Path('uploads')

# Ensure the converted and uploaded directories exist
CONVERTED_FOLDER.mkdir(exist_ok=True)
UPLOADED_FOLDER.mkdir(exist_ok=True)

def save_metadata(file_path, output_json_path):
    """Save file metadata to a JSON file."""
    file_metadata = {
        "file_path": str(file_path.resolve()),  # Absolute path
        "file_name": file_path.name,
        "file_size": file_path.stat().st_size,
        "created_time": time.ctime(file_path.stat().st_ctime),
        "modified_time": time.ctime(file_path.stat().st_mtime)
    }
    with open(output_json_path, 'w') as json_file:
        json.dump(file_metadata, json_file, indent=4)

def convert_pdf_to_docx(pdf_file, output_docx_file):
    """Convert PDF to DOCX using pdf2docx."""
    cv = Converter(str(pdf_file))
    cv.convert(str(output_docx_file), start=0, end=None)
    cv.close()

def extract_text_from_pdf(pdf_path, output_txt_file):
    """Extract plain text from a PDF using PyMuPDF (fitz)."""
    doc = fitz.open(str(pdf_path))
    extracted_text = ""
    for page_num in range(len(doc)):
        page = doc.load_page(page_num)
        text = page.get_text("text")
        extracted_text += text + "\n\n"
    with open(output_txt_file, 'w', encoding='utf-8') as txt_file:
        txt_file.write(extracted_text)

def extract_text_from_pptx(pptx_path, output_txt_file):
    """Extract text from a PowerPoint file."""
    prs = Presentation(str(pptx_path))
    extracted_text = ""
    for slide_num, slide in enumerate(prs.slides):
        slide_header = f"Slide {slide_num + 1}\n{'=' * 10}\n"
        extracted_text += slide_header
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                extracted_text += shape.text + "\n"
        extracted_text += "\n\n"
    with open(output_txt_file, 'w', encoding='utf-8') as txt_file:
        txt_file.write(extracted_text)

def save_pptx_as_docx(pptx_path, output_docx_file):
    """Convert PowerPoint to DOCX format."""
    prs = Presentation(str(pptx_path))
    doc = Document()
    for slide_num, slide in enumerate(prs.slides):
        doc.add_heading(f'Slide {slide_num + 1}', level=1)
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                doc.add_paragraph(shape.text)
    doc.save(output_docx_file)

def process_file(file_path):
    """Process an individual PDF or PPTX file."""
    if file_path.suffix == '.pdf':
        docx_filename = file_path.stem + '.docx'
        docx_path = CONVERTED_FOLDER / docx_filename
        convert_pdf_to_docx(file_path, docx_path)

        txt_filename = file_path.stem + '.txt'
        txt_path = CONVERTED_FOLDER / txt_filename
        extract_text_from_pdf(file_path, txt_path)

        json_filename = file_path.stem + '.json'
        json_path = CONVERTED_FOLDER / json_filename
        save_metadata(file_path, json_path)

    elif file_path.suffix == '.pptx':
        docx_filename = file_path.stem + '.docx'
        docx_path = CONVERTED_FOLDER / docx_filename
        save_pptx_as_docx(file_path, docx_path)

        txt_filename = file_path.stem + '.txt'
        txt_path = CONVERTED_FOLDER / txt_filename
        extract_text_from_pptx(file_path, txt_path)

        json_filename = file_path.stem + '.json'
        json_path = CONVERTED_FOLDER / json_filename
        save_metadata(file_path, json_path)

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """Upload a single file."""
    if not file.filename:
        raise HTTPException(status_code=400, detail="No file uploaded")

    file_path = UPLOADED_FOLDER / file.filename
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    process_file(file_path)
    return {"filename": file.filename}

@app.get("/download_all")
async def download_all():
    """Download all converted files as a zip archive."""
    zip_filename = 'converted_files.zip'
    shutil.make_archive('converted_files', 'zip', CONVERTED_FOLDER)
    return FileResponse(zip_filename, media_type='application/zip', filename=zip_filename)

if __name__ == '__main__':
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)
